from __future__ import annotations

import csv
import json
import math
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/Users/valentinbaumer/Desktop/Probabilistische-Energiedatenvorhersage")
OUT = ROOT / "outputs"
DATA = ROOT / "smard_daten" / "output"
PPTX_PATH = OUT / "Test_ausgefuellt_canva_style.pptx"


NAVY = RGBColor(10, 48, 91)
NAVY_DARK = RGBColor(6, 33, 64)
INK = RGBColor(13, 18, 26)
MUTED = RGBColor(93, 101, 113)
LIGHT = RGBColor(245, 247, 249)
BOX = RGBColor(241, 244, 247)
LINE = RGBColor(178, 188, 199)
ORANGE = RGBColor(255, 111, 61)
BLUE = RGBColor(46, 129, 235)
CYAN = RGBColor(60, 180, 230)
RED = RGBColor(230, 78, 88)
PURPLE = RGBColor(143, 86, 224)
GREEN = RGBColor(38, 134, 88)
WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def rgb(value: RGBColor):
    return value


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor = LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=16, color=INK, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        if bullet:
            p.text = f"• {line}"
    return box


def rect(slide, x, y, w, h, fill=BOX, line=LINE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    if line is None:
        shape.line.fill.background()
    else:
        set_line(shape, line, 0.8)
    return shape


def header(slide, title, subtitle=None):
    add_text(slide, title, 0.7, 0.42, 10.8, 0.48, size=25, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.03))
    set_fill(line, NAVY)
    line.line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 0.7, 1.12, 11.2, 0.28, size=12.5, color=MUTED)


def footer(slide, n):
    add_text(slide, "Probabilistische Energiedatenvorhersage", 0.7, 7.05, 4.4, 0.2, size=7.8, color=MUTED, bold=True)
    add_text(slide, f"{n:02d}", 12.3, 7.05, 0.35, 0.2, size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, accent=BLUE, title_size=13, body_size=10.5):
    rect(slide, x, y, w, h, fill=BOX, line=LINE, radius=True)
    bar = rect(slide, x, y, w, 0.08, fill=accent, line=None)
    add_text(slide, title, x + 0.18, y + 0.18, w - 0.32, 0.25, size=title_size, bold=True, color=accent)
    add_multiline(slide, body if isinstance(body, list) else [body], x + 0.18, y + 0.58, w - 0.35, h - 0.72, size=body_size, color=INK)
    return bar


def image_contain(slide, image_path: Path, x, y, w, h):
    if not image_path.exists():
        rect(slide, x, y, w, h, fill=LIGHT, line=LINE, radius=True)
        add_text(slide, image_path.name, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        return
    rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True)
    with Image.open(image_path) as img:
        iw, ih = img.size
    scale = min((w - 0.18) / iw, (h - 0.18) / ih)
    pic_w = iw * scale
    pic_h = ih * scale
    px = x + (w - pic_w) / 2
    py = y + (h - pic_h) / 2
    slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(pic_w), height=Inches(pic_h))


def metric(slide, x, y, w, h, value, label, note="", color=NAVY):
    rect(slide, x, y, w, h, fill=BOX, line=None, radius=True)
    add_text(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.38, size=23, color=color, bold=True)
    add_text(slide, label, x + 0.18, y + 0.66, w - 0.36, 0.2, size=9.5, bold=True, color=INK)
    if note:
        add_text(slide, note, x + 0.18, y + 0.95, w - 0.36, h - 1.0, size=8.5, color=MUTED)


def read_metrics(prefix):
    path = DATA / f"{prefix}_metriken_04_2026.csv"
    result = {}
    if path.exists():
        with path.open(newline="", encoding="utf-8") as f:
            for row in csv.DictReader(f, delimiter=";"):
                result[row["metric"]] = float(row["value"])
    return result


def fmt_de(value, digits=1):
    if value is None:
        return "-"
    return f"{value:,.{digits}f}".replace(",", "X").replace(".", ",").replace("X", ".")


lstm_metrics = read_metrics("lstm")
rnn_metrics = read_metrics("rnn")


def slide_cover():
    slide = blank_slide()
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    rect(slide, 0.05, 6.67, 13.18, 0.78, fill=NAVY_DARK, line=None)
    add_text(slide, "Probabilistische\nEnergiedaten-\nvorhersage", 0.68, 0.72, 5.8, 1.75, size=30, color=WHITE, bold=True)
    add_text(slide, "Technische Projektvorstellung: Von SMARD-Daten zu Quantilprognosen\nmit LSTM und RNN", 0.7, 3.88, 6.0, 0.55, size=11.5, color=WHITE)
    add_text(slide, "von Jan-Christian Raabe und Valentin Bäumer", 0.7, 4.75, 4.6, 0.22, size=8.8, color=RGBColor(209, 219, 230))
    add_text(slide, "SMARD-Marktdaten · PyTorch · Quantilregression · Streamlit · Docker", 0.68, 6.95, 6.0, 0.22, size=8, color=WHITE, bold=True)
    # Abstract PV icon.
    cx, cy = 9.85, 3.0
    sun = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.28), Inches(cy - 0.28), Inches(0.56), Inches(0.56))
    set_fill(sun, WHITE)
    sun.line.fill.background()
    for angle in range(0, 360, 45):
        x1 = cx + math.cos(math.radians(angle)) * 0.48
        y1 = cy + math.sin(math.radians(angle)) * 0.48
        x2 = cx + math.cos(math.radians(angle)) * 0.78
        y2 = cy + math.sin(math.radians(angle)) * 0.78
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = WHITE
        line.line.width = Pt(2)
    for row in range(2):
        for col in range(6):
            p = rect(slide, 8.65 + col * 0.34 + row * 0.1, 4.45 + row * 0.28, 0.28, 0.18, fill=WHITE, line=None)


def slide_agenda():
    slide = blank_slide()
    header(slide, "Gliederung")
    rows = [
        ("01", "Motivation und Ziel", "Warum probabilistische Prognosen für Photovoltaik sinnvoll sind."),
        ("02", "Gesamtpipeline", "Wie Daten, Modelle, Metriken und Dashboard zusammenarbeiten."),
        ("03", "Verfahren im Detail", "Was, Warum und Wie für Datenaufbereitung, Features, Split, Skalierung, Sequenzen, Quantile, LSTM und RNN."),
        ("04", "Bewertung und Nachweis", "Welche Metriken genutzt werden und wie ein Ergebnisbeispiel aussieht."),
        ("05", "Stärken, Grenzen, Ausblick", "Was das Projekt gut kann, was kritisch bleibt und was der nächste Qualitätssprung wäre."),
    ]
    for i, (num, topic, desc) in enumerate(rows):
        y = 1.28 + i * 0.76
        rect(slide, 0.75, y, 11.55, 0.54, fill=WHITE, line=LINE, radius=True)
        add_text(slide, num, 1.02, y + 0.18, 0.35, 0.15, size=8.5, color=ORANGE, bold=True)
        add_text(slide, topic, 1.7, y + 0.17, 2.25, 0.16, size=8.8, bold=True)
        add_text(slide, desc, 4.55, y + 0.16, 7.15, 0.2, size=8.8, color=INK)
    footer(slide, 2)


def slide_motivation():
    slide = blank_slide()
    header(slide, "Motivation und Zielsetzung", "Volatile PV-Erzeugung erzeugt Risiko, wenn Prognosen nur einen Punktwert liefern.")
    add_text(slide, "Kernproblem", 0.75, 1.55, 2.1, 0.3, size=16, bold=True, color=NAVY)
    add_multiline(
        slide,
        [
            "PV-Erzeugung ist wetter-, tages- und jahreszeitabhängig.",
            "Klassische Punktprognosen zeigen nur einen Wert und blenden Unsicherheit aus.",
            "Bei Abweichungen steigt der Bedarf an kurzfristiger Ausgleichsenergie.",
        ],
        0.75,
        1.95,
        5.3,
        0.98,
        size=12.4,
        color=INK,
    )
    metric(slide, 6.3, 1.55, 1.55, 1.15, "q0.1", "untere Grenze", "80%-Band unten", BLUE)
    metric(slide, 8.15, 1.55, 1.55, 1.15, "q0.5", "Median", "zentrale Prognose", NAVY)
    metric(slide, 10.0, 1.55, 1.55, 1.15, "q0.9", "obere Grenze", "80%-Band oben", GREEN)
    rect(slide, 0.75, 3.35, 11.55, 1.25, fill=BOX, line=None, radius=True)
    add_text(slide, "Ziel des Projekts", 1.0, 3.6, 2.0, 0.25, size=15, bold=True, color=NAVY)
    add_multiline(
        slide,
        [
            "Entwicklung, Implementierung und Evaluierung einer probabilistischen, multivariaten PV-Zeitreihenvorhersage.",
            "Statt einer starren Punktprognose wird ein dynamischer 80%-Sicherheitskorridor berechnet.",
            "Die Umsetzung erfolgt über Quantilregression mit Pinball Loss.",
        ],
        3.1,
        3.42,
        8.0,
        0.9,
        size=12.3,
        color=INK,
    )
    card(
        slide,
        0.75,
        5.1,
        3.45,
        1.25,
        "Modellvergleich",
        ["LSTM und klassisches RNN werden empirisch gegenübergestellt.", "Gesucht ist die robustere kurzfristige PV-Prognose."],
        BLUE,
        body_size=9.2,
    )
    card(
        slide,
        4.55,
        5.1,
        3.45,
        1.25,
        "Transparenz",
        ["Permutation Feature Importance zeigt, welche Einflussfaktoren die Prognose treiben."],
        GREEN,
        body_size=9.2,
    )
    card(
        slide,
        8.35,
        5.1,
        3.45,
        1.25,
        "Validierung",
        ["Das Unsicherheitsband wird gegen historische Realität geprüft:", "breiter bei hoher Unsicherheit, enger bei stabileren Verläufen."],
        ORANGE,
        body_size=9.2,
    )
    footer(slide, 3)


def slide_pipeline():
    slide = blank_slide()
    header(slide, "Gesamtpipeline", "Die Komponenten greifen als reproduzierbare Pipeline ineinander.")
    labels = [
        ("Upload", "SMARD-CSV\noptional\nWetter-CSV", ORANGE),
        ("EDA", "Bereinigung\nTagesprofile\nKorrelationen", ORANGE),
        ("Features", "Zeitzyklen\nPV-Historie\nWind, Erdgas", RED),
        ("Split & Scale", "chronologisch\nScaler nur Train", PURPLE),
        ("Modelle", "LSTM und RNN\nq0.1, q0.5, q0.9", BLUE),
        ("Evaluation", "PICP, RMSE\nWinkler,\nPinball", CYAN),
    ]
    shapes = []
    for i, (title, body, color) in enumerate(labels):
        x = 0.7 + i * 2.05
        rect(slide, x, 2.25, 1.6, 1.25, fill=BOX, line=LINE, radius=True)
        rect(slide, x, 2.25, 1.6, 0.08, fill=color, line=None)
        add_text(slide, title, x + 0.13, 2.43, 1.25, 0.2, size=8.5, bold=True, color=color)
        add_text(slide, body, x + 0.13, 2.82, 1.22, 0.5, size=7.2, color=INK)
        shapes.append((x, 2.25))
        if i < len(labels) - 1:
            conn = slide.shapes.add_connector(1, Inches(x + 1.6), Inches(2.9), Inches(x + 2.0), Inches(2.9))
            conn.line.color.rgb = LINE
            conn.line.width = Pt(1.2)
    # Output loop to dashboard.
    line = slide.shapes.add_connector(2, Inches(11.55), Inches(3.52), Inches(6.7), Inches(4.25))
    line.line.color.rgb = LINE
    line.line.width = Pt(1.1)
    rect(slide, 4.25, 4.9, 6.1, 0.82, fill=NAVY, line=None, radius=True)
    add_text(slide, "Streamlit-Dashboard", 4.65, 5.18, 2.3, 0.2, size=9, color=WHITE, bold=True)
    add_text(slide, "macht Metriken, Plots, Protokolle und regelbasierte Interpretation direkt nutzbar", 7.05, 5.05, 2.75, 0.42, size=8, color=WHITE)
    footer(slide, 4)


def slide_eda():
    slide = blank_slide()
    header(slide, "Explorative Datenanalyse (EDA)", "Die Datenaufbereitung sorgt dafür, dass aus CSV-Zeilen ein lernbarer Zeitindex wird.")
    card(slide, 0.75, 1.75, 3.4, 2.2, "Was?", "Die SMARD-Datei wird eingelesen, deutsche Zahlenformate werden konvertiert und Zeitstempel werden als sortierter DatetimeIndex verarbeitet.", CYAN)
    card(slide, 4.55, 1.75, 3.4, 2.2, "Warum?", "Neuronale Modelle können nur mit konsistenten numerischen Arrays arbeiten. Fehlerhafte Zeitstempel würden Training und Evaluation verfälschen.", BLUE)
    card(slide, 8.35, 1.75, 3.4, 2.2, "Wie?", "Flexible CSV-Erkennung, Umwandlung bspw. von 4.167,28 zu 4167.28, Sortierung nach Zeit und Entfernen bzw. Interpolieren fehlender Werte.", PURPLE)
    rect(slide, 0.75, 4.38, 11.0, 0.72, fill=WHITE, line=LINE, radius=True)
    add_text(slide, "Optional kann eine Wetter-CSV ergänzt werden. Numerische Wetterspalten werden mit `weather_` präfixiert und zeitlich auf die SMARD-Reihe gemergt.", 0.98, 4.62, 10.5, 0.25, size=9.5, color=INK, align=PP_ALIGN.CENTER)
    image_contain(slide, DATA / "eda_tagesprofil_04_2026.png", 0.75, 5.45, 5.0, 1.45)
    add_multiline(
        slide,
        [
            "Das Tagesprofil zeigt, warum Zeitfeatures sinnvoll sind:",
            "Photovoltaik folgt einem klaren Tagesrhythmus.",
        ],
        6.1,
        5.72,
        5.0,
        0.65,
        size=12.6,
        color=INK,
    )
    footer(slide, 5)


def slide_features():
    slide = blank_slide()
    header(slide, "Feature Engineering", "Zeitliche Muster und Kontextinformationen werden explizit lernbar gemacht.")
    image_contain(slide, DATA / "eda_korrelations_heatmap_04_2026.png", 0.75, 1.55, 4.6, 3.15)
    card(slide, 5.75, 1.55, 5.95, 0.9, "Was?", "Aus Rohdaten werden Modellfeatures: PV-Historie, Wind Onshore, Erdgas, Stunde, Wochentag und Wochenende.", BLUE, body_size=9)
    card(slide, 5.75, 2.75, 5.95, 0.9, "Warum?", "PV folgt zyklischen Mustern. Sinus/Cosinus kodieren Uhrzeit und Wochentag so, dass Übergänge wie 23 Uhr zu 0 Uhr sinnvoll abgebildet werden.", GREEN, body_size=9)
    card(slide, 5.75, 3.95, 5.95, 0.9, "Wie?", "`hour_sin`, `hour_cos`, `weekday_sin`, `weekday_cos`, `is_weekend`; optionale Wetterspalten werden automatisch ergänzt.", ORANGE, body_size=9)
    rect(slide, 0.75, 5.35, 11.0, 0.92, fill=BOX, line=None, radius=True)
    add_text(slide, "Begründung der Feature-Auswahl", 1.0, 5.62, 2.3, 0.2, size=12.5, bold=True, color=NAVY)
    add_multiline(
        slide,
        [
            "Die Features verbinden historische PV-Erzeugung,",
            "zeitliche Struktur und weitere Energiesystemgrößen.",
        ],
        3.55,
        5.48,
        7.2,
        0.55,
        size=11.8,
        color=INK,
    )
    footer(slide, 6)


def slide_split_scaling():
    slide = blank_slide()
    header(slide, "Chronologischer Split und Skalierung", "Die wichtigste methodische Absicherung gegen Data Leakage.")
    card(slide, 0.75, 1.55, 3.4, 2.55, "Was?", "Die Zeitreihe wird chronologisch in Training, Validierung und Test aufgeteilt. Danach werden Features und Zielvariable skaliert.", BLUE)
    card(slide, 4.55, 1.55, 3.4, 2.55, "Warum?", "Ein zufälliger Split oder ein Scaler auf allen Daten würde Informationen aus dem Testzeitraum ins Training tragen.", GREEN)
    card(slide, 8.35, 1.55, 3.4, 2.55, "Wie?", "`MinMaxScaler` wird nur auf Trainingsdaten gefittet. Validierung und Test werden danach mit demselben Scaler transformiert.", ORANGE)
    rect(slide, 2.2, 4.9, 8.9, 0.75, fill=NAVY_DARK, line=None, radius=True)
    add_text(slide, "fit scaler: train only    →    transform train + validation + test", 2.55, 5.17, 8.1, 0.18, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "So bewerten die Testmetriken echte Zukunftsdaten relativ zum Trainingszeitraum.", 0.75, 6.08, 11.1, 0.24, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 7)


def slide_sequences():
    slide = blank_slide()
    header(slide, "Sequenzbau", "Aus Einzelzeitpunkten wird ein echtes Zeitreihenproblem.")
    add_text(slide, "24 Stunden Vergangenheit", 0.85, 1.65, 3.5, 0.28, size=16, bold=True, color=NAVY)
    for i in range(24):
        x = 0.85 + (i % 12) * 0.25
        y = 2.2 + (i // 12) * 0.38
        rect(slide, x, y, 0.18, 0.18, fill=BLUE if i < 12 else CYAN, line=None)
    arrow = slide.shapes.add_connector(1, Inches(4.2), Inches(2.45), Inches(5.5), Inches(2.45))
    arrow.line.color.rgb = NAVY
    arrow.line.width = Pt(2)
    add_text(slide, "Modell", 5.75, 2.22, 1.0, 0.25, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, 5.55, 2.0, 1.45, 0.86, fill=NAVY, line=None, radius=True)
    arrow2 = slide.shapes.add_connector(1, Inches(7.0), Inches(2.45), Inches(8.2), Inches(2.45))
    arrow2.line.color.rgb = NAVY
    arrow2.line.width = Pt(2)
    rect(slide, 8.35, 2.02, 2.5, 0.82, fill=BOX, line=LINE, radius=True)
    add_text(slide, "nächster Zielwert", 8.7, 2.27, 1.8, 0.2, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    card(slide, 0.75, 4.05, 3.4, 1.55, "Was?", "Jede Vorhersage nutzt ein Fenster vergangener Werte statt nur einen Zeitpunkt.", BLUE)
    card(slide, 4.55, 4.05, 3.4, 1.55, "Warum?", "Tagesanstieg, Peak und Abendabfall sind zeitliche Muster, die Sequenzmodelle lernen können.", GREEN)
    card(slide, 8.35, 4.05, 3.4, 1.55, "Wie?", "Die Auflösung bestimmt die Fensterlänge. Bei Stundendaten entspricht 24h einem Fenster von 24 Schritten.", ORANGE)
    footer(slide, 8)


def slide_quantiles():
    slide = blank_slide()
    header(slide, "Quantilprognose", "Die reine Punktprognose wird durch Unsicherheitsinformation ersetzt.")
    rect(slide, 0.85, 1.55, 4.85, 3.4, fill=BOX, line=None, radius=True)
    add_text(slide, "q0.9", 1.25, 2.0, 0.8, 0.25, size=20, bold=True, color=GREEN)
    add_text(slide, "obere Grenze", 2.05, 2.03, 1.6, 0.2, size=12, color=INK)
    add_text(slide, "q0.5", 1.25, 3.0, 0.8, 0.25, size=20, bold=True, color=NAVY)
    add_text(slide, "Median", 2.05, 3.03, 1.6, 0.2, size=12, color=INK)
    add_text(slide, "q0.1", 1.25, 4.0, 0.8, 0.25, size=20, bold=True, color=BLUE)
    add_text(slide, "untere Grenze", 2.05, 4.03, 1.6, 0.2, size=12, color=INK)
    for y, color in [(2.14, GREEN), (3.14, NAVY), (4.14, BLUE)]:
        line = slide.shapes.add_connector(1, Inches(3.65), Inches(y), Inches(5.0), Inches(y))
        line.line.color.rgb = color
        line.line.width = Pt(2.2)
    card(slide, 6.15, 1.55, 5.35, 1.0, "Was?", "Das Modell gibt drei Werte aus: q0.1, q0.5 und q0.9.", BLUE)
    card(slide, 6.15, 2.85, 5.35, 1.0, "Warum?", "Ein einzelner Wert unterschätzt die Unsicherheit volatiler PV-Erzeugung.", GREEN)
    card(slide, 6.15, 4.15, 5.35, 1.0, "Wie?", "Die letzte lineare Schicht hat drei Outputs. q0.1 bis q0.9 bilden ein nominelles 80%-Intervall.", ORANGE)
    rect(slide, 0.85, 5.65, 10.65, 0.52, fill=WHITE, line=LINE, radius=True)
    add_text(slide, "Die 80%-Abdeckung ist bewusst gewählt: gut erklärbar, breit genug für Unsicherheit, aber nicht so breit wie ein 95%-Intervall.", 1.05, 5.83, 10.2, 0.2, size=10.8, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 9)


def slide_pinball():
    slide = blank_slide()
    header(slide, "Pinball Loss und Crossing-Penalty", "Quantile brauchen eine andere Verlustfunktion als klassische Punktprognosen.")
    card(slide, 0.75, 1.55, 3.4, 2.35, "Was?", "Pinball Loss ist eine asymmetrische Verlustfunktion für Quantile. Unter- und Überschätzung werden je nach Quantil unterschiedlich bestraft.", BLUE)
    card(slide, 4.55, 1.55, 3.4, 2.35, "Warum?", "MSE und MAE trainieren zentrale Punktwerte. Für q0.1 und q0.9 brauchen wir gezielte untere und obere Grenzen.", GREEN)
    card(slide, 8.35, 1.55, 3.4, 2.35, "Wie?", "Die Quantil-Losses werden summiert. Zusätzlich bestraft eine Crossing-Penalty falsche Quantilreihenfolgen.", ORANGE)
    rect(slide, 2.0, 4.72, 9.4, 0.86, fill=NAVY_DARK, line=None, radius=True)
    add_text(slide, "training_loss = pinball_loss(q0.1, q0.5, q0.9) + 0.2 · crossing_penalty", 2.4, 5.03, 8.6, 0.22, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Ziel: fachlich konsistente Prognosen mit q0.1 ≤ q0.5 ≤ q0.9.", 0.75, 6.2, 11.0, 0.25, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 10)


def slide_models():
    slide = blank_slide()
    header(slide, "LSTM und RNN", "Zwei Sequenzmodelle mit identischer Trainings- und Evaluationslogik.")
    rect(slide, 0.85, 1.55, 5.2, 3.95, fill=BOX, line=LINE, radius=True)
    rect(slide, 7.0, 1.55, 5.2, 3.95, fill=BOX, line=LINE, radius=True)
    add_text(slide, "LSTM", 1.15, 1.9, 1.5, 0.35, size=22, bold=True, color=BLUE)
    add_text(slide, "Was?", 1.15, 2.55, 0.9, 0.22, size=11, bold=True, color=BLUE)
    add_text(slide, "Rekurrentes Netz mit Gates zum Speichern und Vergessen von Information.", 2.0, 2.52, 3.45, 0.38, size=11)
    add_text(slide, "Warum?", 1.15, 3.25, 0.9, 0.22, size=11, bold=True, color=GREEN)
    add_text(slide, "Geeignet, wenn längere zeitliche Abhängigkeiten relevant sind.", 2.0, 3.22, 3.45, 0.38, size=11)
    add_text(slide, "Wie?", 1.15, 3.95, 0.9, 0.22, size=11, bold=True, color=ORANGE)
    add_text(slide, "2 Layer, Hidden Size 128, Dropout 0.2, lineare Ausgabe auf drei Quantile.", 2.0, 3.92, 3.45, 0.42, size=11)
    add_text(slide, "RNN", 7.3, 1.9, 1.5, 0.35, size=22, bold=True, color=BLUE)
    add_text(slide, "Was?", 7.3, 2.55, 0.9, 0.22, size=11, bold=True, color=BLUE)
    add_text(slide, "Einfaches rekurrentes Netz mit verstecktem Zustand pro Zeitschritt.", 8.15, 2.52, 3.45, 0.38, size=11)
    add_text(slide, "Warum?", 7.3, 3.25, 0.9, 0.22, size=11, bold=True, color=GREEN)
    add_text(slide, "Verständliche Baseline: zeigt, ob die komplexere LSTM-Struktur wirklich nötig ist.", 8.15, 3.22, 3.45, 0.44, size=11)
    add_text(slide, "Wie?", 7.3, 3.95, 0.9, 0.22, size=11, bold=True, color=ORANGE)
    add_text(slide, "Gleiche Featurebasis, gleiche Quantile, gleiche Evaluation; nur Kernarchitektur unterscheidet sich.", 8.15, 3.92, 3.45, 0.46, size=11)
    rect(slide, 2.2, 5.95, 8.8, 0.58, fill=NAVY, line=None, radius=True)
    add_text(slide, "Fairer Vergleich: Training, Split, Loss, Metriken und Plots liegen in gemeinsamen Hilfsfunktionen.", 2.5, 6.16, 8.25, 0.18, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 11)


def slide_training():
    slide = blank_slide()
    header(slide, "Training und Reproduzierbarkeit", "Der Modelllauf soll nachvollziehbar und fair vergleichbar sein.")
    items = [
        ("Seed 42", "Python, NumPy und PyTorch werden initialisiert.", BLUE),
        ("DataLoader", "Training läuft in Batches mit reproduzierbarem Shuffle.", CYAN),
        ("Early Stopping", "Der beste Checkpoint wird über den Validierungsloss gewählt.", GREEN),
        ("Gemeinsame Pipeline", "LSTM und RNN unterscheiden sich nur in der Modellklasse.", ORANGE),
        ("Artefakte", "Modelle, Scaler, Metriken, Plots und Metadaten werden gespeichert.", PURPLE),
    ]
    for i, (title, body, color) in enumerate(items):
        x = 0.8 + (i % 3) * 3.9
        y = 1.55 + (i // 3) * 1.9
        card(slide, x, y, 3.3, 1.35, title, body, color)
    rect(slide, 0.8, 5.95, 11.0, 0.6, fill=WHITE, line=LINE, radius=True)
    add_text(slide, "Begründung: Wenn beide Modelle gleich behandelt werden, ist der Vergleich methodisch deutlich sauberer.", 1.0, 6.17, 10.6, 0.18, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 12)


def slide_evaluation():
    slide = blank_slide()
    header(slide, "Evaluation", "Ein probabilistisches Modell muss Fehler und Unsicherheit gemeinsam bestehen.")
    metrics = [
        ("PICP", "Wie viele echte Werte liegen im 80%-Intervall?"),
        ("Breite", "Wie informativ oder breit ist das Unsicherheitsband?"),
        ("Median-RMSE", "Wie stark weicht die q0.5-Prognose ab?"),
        ("Pinball Loss", "Wie gut sind die einzelnen Quantile?"),
        ("Winkler Score", "Wie ist die Balance aus Breite und Fehldeckung?"),
        ("Crossing Rate", "Bleibt q0.1 ≤ q0.5 ≤ q0.9 erhalten?"),
    ]
    for i, (name, desc) in enumerate(metrics):
        x = 0.8 + (i % 2) * 5.65
        y = 1.55 + (i // 2) * 1.1
        rect(slide, x, y, 5.1, 0.74, fill=BOX, line=LINE, radius=True)
        add_text(slide, name, x + 0.25, y + 0.24, 1.45, 0.18, size=12.5, bold=True, color=GREEN if i == 0 else NAVY)
        add_text(slide, desc, x + 1.85, y + 0.18, 2.95, 0.26, size=10.2)
    rect(slide, 0.8, 5.35, 11.0, 0.74, fill=WHITE, line=LINE, radius=True)
    add_text(slide, "Wichtig: Ein niedriger RMSE allein reicht nicht. Ein Intervall muss auch kalibriert und nicht unnötig breit sein.", 1.05, 5.62, 10.5, 0.22, size=12.5, color=MUTED, align=PP_ALIGN.CENTER)
    footer(slide, 13)


def slide_dashboard():
    slide = blank_slide()
    header(slide, "Streamlit-Dashboard", "Die technischen Artefakte werden für Nutzerinnen und Nutzer nachvollziehbar aufbereitet.")
    card(slide, 0.8, 1.55, 3.5, 2.15, "Was?", "Upload, Pipeline-Start, Ergebnisse, Diagramme, Metriken, Dateien, Protokoll und Impressum liegen in einer lokalen Oberfläche.", BLUE)
    card(slide, 4.9, 1.55, 3.5, 2.15, "Warum?", "Das Projekt soll nicht nur Skripte liefern, sondern einen nachvollziehbaren Workflow für Außenstehende.", GREEN)
    card(slide, 9.0, 1.55, 3.0, 2.15, "Wie?", "Nach der Pipeline wird ein Dashboard-Bundle erzeugt, das kompakte Metriken und Vorschaubilder lädt.", ORANGE)
    rect(slide, 1.0, 4.4, 10.65, 0.9, fill=NAVY_DARK, line=None, radius=True)
    add_text(slide, "Regelbasierte Interpretation", 1.35, 4.68, 3.0, 0.25, size=14.5, bold=True, color=WHITE)
    add_text(slide, "PICP, Crossing, RMSE, Winkler Score und Modellvergleich werden mit festen, transparenten Regeln erklärt.", 4.35, 4.62, 6.6, 0.32, size=11.5, color=WHITE)
    footer(slide, 14)


def slide_results():
    slide = blank_slide()
    header(slide, "Ergebnisbeispiel als Funktionsnachweis", "Die Ergebnisse zeigen, dass die Pipeline von Upload bis Evaluation durchläuft.")
    image_contain(slide, DATA / "rnn_vorhersage_intervall_04_2026.png", 0.75, 1.45, 7.4, 4.0)
    metric(slide, 8.55, 1.55, 1.55, 1.1, f"{fmt_de(rnn_metrics.get('PICP_80_percent'), 1)}%", "RNN PICP", "nahe, eher konservativ", GREEN)
    metric(slide, 10.35, 1.55, 1.55, 1.1, f"{fmt_de(rnn_metrics.get('Quantile_Crossing_Rate_percent'), 1)}%", "Crossing", "Output stabil", GREEN)
    metric(slide, 8.55, 3.0, 1.55, 1.1, fmt_de(rnn_metrics.get("Median_RMSE_MWh"), 0), "RMSE", "MWh", NAVY)
    metric(slide, 10.35, 3.0, 1.55, 1.1, fmt_de(rnn_metrics.get("Winkler_Score_80_MWh"), 0), "Winkler", "MWh", NAVY)
    add_text(slide, "Einordnung", 8.55, 4.55, 1.9, 0.25, size=14, bold=True, color=NAVY)
    add_multiline(
        slide,
        [
            "Das Beispiel ist kein endgültiger Produktivnachweis.",
            "Es zeigt, dass die probabilistische Pipeline technisch funktioniert.",
        ],
        8.55,
        4.92,
        3.35,
        0.75,
        size=10.8,
        color=INK,
    )
    footer(slide, 15)


def slide_importance():
    slide = blank_slide()
    header(slide, "Interpretierbarkeit", "Feature Importance zeigt, welche Eingangsinformation das Modell nutzt.")
    image_contain(slide, DATA / "rnn_feature_importance_04_2026.png", 0.75, 1.45, 6.4, 4.2)
    card(slide, 7.65, 1.55, 4.0, 1.0, "Wichtigster Treiber", "Die eigene PV-Historie liefert im Beispiel den größten Beitrag.", GREEN)
    card(slide, 7.65, 2.85, 4.0, 1.0, "Zeitliche Struktur", "Zeitfeatures erklären den Tagesrhythmus und ergänzen die Historie.", BLUE)
    card(slide, 7.65, 4.15, 4.0, 1.0, "Kritischer Punkt", "Ohne echte Wetterfeatures sieht das Modell keine konkrete Wetterlage.", RED)
    footer(slide, 16)


def slide_strengths():
    slide = blank_slide()
    header(slide, "Was das Projekt bereits gut kann", "Stärken, die in der Präsentation ruhig betont werden können.")
    strengths = [
        ("Saubere Methodik", "Chronologischer Split und Scaler nur auf Training reduzieren Data Leakage."),
        ("Probabilistische Sicht", "Die Prognose zeigt Unsicherheit statt nur einen Zielwert."),
        ("Vergleichbarkeit", "LSTM und RNN teilen Training, Sequenzbau und Evaluation."),
        ("Nachvollziehbarkeit", "Metriken, Plots, Protokolle und Dashboard machen Ergebnisse prüfbar."),
        ("Reproduzierbarkeit", "Seeds, Dependencies, Docker Compose und Artefaktstruktur helfen bei Wiederholung."),
        ("Erweiterbarkeit", "Wetterfeatures sind bereits vorgesehen und fachlich sinnvoll ergänzbar."),
    ]
    for i, (title, body) in enumerate(strengths):
        x = 0.8 + (i % 2) * 5.55
        y = 1.45 + (i // 2) * 1.25
        rect(slide, x, y, 5.0, 0.82, fill=BOX, line=None, radius=True)
        add_text(slide, title, x + 0.22, y + 0.25, 1.8, 0.18, size=10.5, bold=True, color=GREEN)
        add_text(slide, body, x + 2.05, y + 0.17, 2.65, 0.34, size=9.3, color=INK)
    footer(slide, 17)


def slide_limits():
    slide = blank_slide()
    header(slide, "Was man kritisch sehen muss", "Grenzen und nächste Entwicklungsschritte.")
    limits = [
        ("Keine echte Wetterprognose", "Ohne Globalstrahlung, Bewölkung und Temperatur bleibt PV fachlich begrenzt.", RED),
        ("Kleine Datenbasis", "Ein kurzer Zeitraum reicht für einen Prototyp, aber nicht für robuste saisonale Aussagen.", NAVY),
        ("Kein operativer Einsatz", "Das System ist ein Demonstrations- und Forschungsprojekt.", NAVY),
        ("Modellvergleich ist laufabhängig", "Ein einzelner Lauf ist kein allgemeiner Architekturbeweis.", NAVY),
        ("Unsicherheit muss kalibriert sein", "Zu breite Intervalle decken viel ab, sind aber wenig informativ.", NAVY),
        ("Exogene Ereignisse fehlen", "Feiertage, Marktverhalten oder Anlagenänderungen werden nicht explizit modelliert.", NAVY),
    ]
    for i, (title, body, color) in enumerate(limits):
        x = 0.8 + (i % 2) * 5.55
        y = 1.45 + (i // 2) * 1.25
        rect(slide, x, y, 5.0, 0.82, fill=RGBColor(255, 241, 237) if i == 0 else BOX, line=None, radius=True)
        add_text(slide, title, x + 0.22, y + 0.2, 1.8, 0.26, size=9.8, bold=True, color=color)
        add_text(slide, body, x + 2.05, y + 0.15, 2.7, 0.4, size=9.2, color=INK)
    footer(slide, 18)


def slide_thanks():
    slide = blank_slide()
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, "Vielen Dank für eure Aufmerksamkeit!", 0.9, 0.8, 11.5, 0.55, size=25, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # PV icon, larger than cover.
    cx, cy = 6.65, 3.25
    sun = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx - 0.32), Inches(cy - 0.32), Inches(0.64), Inches(0.64))
    set_fill(sun, WHITE)
    sun.line.fill.background()
    for angle in range(0, 360, 45):
        x1 = cx + math.cos(math.radians(angle)) * 0.55
        y1 = cy + math.sin(math.radians(angle)) * 0.55
        x2 = cx + math.cos(math.radians(angle)) * 0.9
        y2 = cy + math.sin(math.radians(angle)) * 0.9
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = WHITE
        line.line.width = Pt(2)
    for row in range(3):
        for col in range(7):
            p = rect(slide, 5.1 + col * 0.33 + row * 0.11, 4.55 + row * 0.25, 0.26, 0.16, fill=WHITE, line=None)
    add_text(slide, "Fragen?", 0.9, 5.9, 11.5, 0.4, size=23, color=WHITE, align=PP_ALIGN.CENTER)


for builder in [
    slide_cover,
    slide_agenda,
    slide_motivation,
    slide_pipeline,
    slide_eda,
    slide_features,
    slide_split_scaling,
    slide_sequences,
    slide_quantiles,
    slide_pinball,
    slide_models,
    slide_training,
    slide_evaluation,
    slide_dashboard,
    slide_results,
    slide_importance,
    slide_strengths,
    slide_limits,
    slide_thanks,
]:
    builder()


OUT.mkdir(exist_ok=True)
prs.save(PPTX_PATH)
print(PPTX_PATH)
